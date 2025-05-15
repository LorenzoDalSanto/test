from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR

def add_textbox(slide, text, left, top, width, height, font_size, font_name, bold, italic, color_rgb, centered=True):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame

    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True

    tf.clear()
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT

    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = font_name
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = RGBColor(*color_rgb)

# In utils.py

# ... (altre importazioni e funzioni) ...

def add_horizontal_line(slide, left, top, width, thickness=Pt(1), color=(150, 150, 150)):
    """
    Aggiunge una linea orizzontale reale (vettoriale) nello slide.
    Nessuna ombra, nessun bordo inutile.
    """
    # Crea connettore (linea vera)
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        left, top,
        left + width, top
    )

    line.line.width = thickness
    line.line.fill.solid()
    line.line.fill.fore_color.rgb = RGBColor(*color)

    # Rimuove ombre e glow (anche se non sempre presenti)
    try:
        line.shadow.inherit = False
        line.line.format.shadow.inherit = False
        line.line.format.shadow.blur_radius = 0
        line.line.format.glow.radius = 0
        line.line.format.soft_edge.radius = 0
    except:
        pass  # fallback in caso l'API non esponga qualcosa

    return line

# ... (resto di utils.py) ...


def add_competitor_row(slide, data, y_offset):

    # SOCIETÀ
    add_textbox(
        slide, text=data["societa"],
        left=435600, top=1360800+y_offset, width=1076400, height=252000,
        font_size=10.5, font_name="Calibri", bold=False, italic=False, color_rgb=(0, 0, 0)
    )

    # SEDE
    add_textbox(
        slide, text=data["sede"],
        left=3078000, top=1270800+y_offset , width=1076400, height=410400,
        font_size=10.5, font_name="Calibri", bold=False, italic=False, color_rgb=(0, 0, 0)
    )

    # DESCRIZIONE
    add_textbox(
        slide, text=data["descrizione"],
        left=4233600, top=1191600+y_offset, width=3528000, height=252000, centered=False,
        font_size=10.5, font_name="Calibri", bold=False, italic=False, color_rgb=(0, 0, 0)
    )

    # VDP
    add_textbox(
        slide, text="".join(["VDP\n", data["vdp"]]),
        left=7920000, top=1310400+y_offset, width=666000, height=363600,
        font_size=9, font_name="Calibri", bold=False, italic=False, color_rgb=(0, 0, 0)
    )

    # EBITDA
    add_textbox(
        slide, text="EBITDA\n"+data["ebitda"],
        left=8575200, top=1310400+y_offset, width=666000, height=363600,
        font_size=9, font_name="Calibri", bold=False, italic=False, color_rgb=(0, 0, 0)
    )

    # EBITDA %
    add_textbox(
        slide, text="EBITDA %\n"+data["ebitda_percent"],
        left=9230400, top=1310400+y_offset, width=666000, height=363600,
        font_size=9, font_name="Calibri", bold=False, italic=False, color_rgb=(0, 0, 0)
    )

    # N DIP
    add_textbox(
        slide, text="N. Dip.\n"+data["n_dip"],
        left=9885600, top=1310400+y_offset, width=666000, height=363600,
        font_size=9, font_name="Calibri", bold=False, italic=False, color_rgb=(0, 0, 0)
    )
    # PFN
    add_textbox(
        slide, text="PFN\n"+data["pfn"],
        left=10540800, top=1310400+y_offset, width=666000, height=363600,
        font_size=9, font_name="Calibri", bold=False, italic=False, color_rgb=(0, 0, 0)
    )

    # FCF
    add_textbox(
        slide, text="FCF\n"+data["fcf"],
        left=11196000, top=1310400+y_offset, width=666000, height=363600,
        font_size=9, font_name="Calibri", bold=False, italic=False, color_rgb=(0, 0, 0)
    )


def duplicate_slide(prs, slide_index):
    """
    Duplica uno slide (con tutte le forme) in una nuova posizione alla fine della presentazione.
    Restituisce il nuovo slide.
    """
    source = prs.slides[slide_index]
    layout = source.slide_layout
    new_slide = prs.slides.add_slide(layout)

    for shape in source.shapes:
        el = shape.element
        new_el = el.clone()
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return new_slide


def add_shapes(slide):

    add_textbox(
        slide, text="Società",
        left=388800, top=813600, width=1119600, height=302400,
        font_size=14, font_name="Calibri", bold=True, italic=False, color_rgb=(0, 0, 0)
    )

    add_textbox(
        slide, text="Logo",
        left=1688400, top=813600, width=1119600, height=302400,
        font_size=14, font_name="Calibri", bold=True, italic=False, color_rgb=(0, 0, 0)
    )

    add_textbox(
        slide, text="Sede",
        left=2984400, top=813600, width=1119600, height=302400,
        font_size=14, font_name="Calibri", bold=True, italic=False, color_rgb=(0, 0, 0)
    )

    add_textbox(
        slide, text="Descrizione dell'attività svolta",
        left=4719600, top=813600, width=2548800, height=302400,
        font_size=14, font_name="Calibri", bold=True, italic=False, color_rgb=(0, 0, 0)
    )

    add_textbox(
        slide, text="Key Financials (€/k)",
        left=8766000, top=813600, width=2185200, height=302400,
        font_size=14, font_name="Calibri", bold=True, italic=False, color_rgb=(0, 0, 0)
    )

    add_textbox(
        slide, text="Competitors Analysis",
        left=342000, top=198000, width=10656000, height=363600, centered=False,
        font_size=24, font_name="Calibri", bold=False, italic=False, color_rgb=(0, 0, 0)
    )
    
    add_horizontal_line(
        slide,
        left=435600, top=1893600, width=11325600,
    )
    
    add_horizontal_line(
        slide,
        left=435600, top=2653200, width=11325600,
    )
    
    add_horizontal_line(
        slide,
        left=435600, top=3412800, width=11325600,
    )
    
    add_horizontal_line(
        slide,
        left=435600, top=4183200, width=11325600,
    )
    
    add_horizontal_line(
        slide,
        left=435600, top=4975200, width=11325600,
    )
    
    add_horizontal_line(
        slide,
        left=435600, top=5774400, width=11325600,
    )

    add_horizontal_line(
        slide,
        left=345600, top=1105200, width=1206000,
        thickness=Pt(2), color=(0, 0, 0)
    )

    add_horizontal_line(
        slide,
        left=1645200, top=1105200, width=1206000,
        thickness=Pt(2), color=(0, 0, 0)
    )

    add_horizontal_line(
        slide,
        left=2944800, top=1105200, width=1206000,
        thickness=Pt(2), color=(0, 0, 0)
    )

    add_horizontal_line(
        slide,
        left=4230000, top=1105200, width=3528000,
        thickness=Pt(2), color=(0, 0, 0)
    )

    add_horizontal_line(
        slide,
        left=7920000, top=1105200, width=3938400,
        thickness=Pt(2), color=(0, 0, 0)
    )
    

