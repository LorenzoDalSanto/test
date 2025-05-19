import pandas as pd
from pptx import Presentation
import utils
import os   # Per interagire con il file system (es. nomi file, percorsi)
import glob # Per trovare file che corrispondono a un pattern (es. *.xlsx)

# --- 1. Definizione percorso e chiavi di estrazione ---
aziende_target_folder = "AziendeTarget"  # Nome della sottocartella
excel_files_pattern = os.path.join(aziende_target_folder, "*.xlsx")

# Chiavi da cercare nei file Excel (originali, come fornite dall'utente)
# e la loro mappatura ai nomi standard usati in add_competitor_row
# Le chiavi del dizionario esterno sono le etichette da cercare (normalizzate in minuscolo)
# I valori sono i nomi standard delle colonne per il DataFrame finale
chiavi_excel_mapping = {
    "descrizione attività italiano": "descrizione", # Modificato per matchare la logica precedente
    "dipendenti": "n_dip",
    "ebitda": "ebitda",
    "ebitda/vendite (%)": "ebitda_percent",
    "posizione finanziaria netta": "pfn",
    "a. tot. val. della produzione": "vdp", # Assumendo che questa sia l'etichetta per Valore Produzione
    "fcf": "fcf",
    "indirizzo sede legale": "sede" # Aggiunta per la sede, da verificare se presente con questa etichetta
}
# Converti le chiavi di ricerca (etichette Excel) in minuscolo per confronto case-insensitive
chiavi_excel_target_lower = {k.lower().strip(): v for k, v in chiavi_excel_mapping.items()}

# Lista di tutte le chiavi standard che ci aspettiamo nel dizionario finale per ogni azienda
# Questo aiuta a garantire che ogni dizionario abbia la stessa struttura.
standard_keys_final = ["societa", "sede", "descrizione", "vdp", "ebitda", "ebitda_percent", "n_dip", "pfn", "fcf"]


# --- 2. Estrazione dati da tutti i file Excel ---
lista_dati_aziende = []

excel_file_paths = glob.glob(excel_files_pattern) # Trova tutti i file .xlsx nella cartella

if not excel_file_paths:
    print(f"ATTENZIONE: Nessun file .xlsx trovato nella cartella '{aziende_target_folder}'.")
    print("Assicurati che la cartella esista e contenga i file Excel.")
    exit()

print(f"Trovati i seguenti file Excel da processare: {excel_file_paths}")

for file_path in excel_file_paths:
    print(f"\n--- Processando file: {file_path} ---")
    try:
        # Leggi l'intero foglio senza assumere una riga di intestazione specifica
        # Questo permette di cercare le etichette in qualsiasi riga.
        df_excel = pd.read_excel(file_path, header=None) 
        
        dati_singola_azienda = {} # Dizionario per i dati del file corrente

        # Estrai nome società dal nome del file (rimuovendo l'estensione .xlsx)
        nome_societa = os.path.splitext(os.path.basename(file_path))[0]
        dati_singola_azienda["societa"] = nome_societa
        print(f"  Società (dal nome file): {nome_societa}")

        # Itera sulle righe del DataFrame letto dal file Excel
        for index, row in df_excel.iterrows():
            # Cerca le etichette nelle prime due colonne (0 e 1)
            for col_idx in range(min(2, len(row))): # len(row) per evitare IndexError se ci sono meno di 2 colonne
                cella_contenuto = row.iloc[col_idx]
                
                if isinstance(cella_contenuto, str):
                    etichetta_cella_lower = cella_contenuto.strip().lower()
                    
                    # Se l'etichetta della cella è una di quelle che cerchiamo
                    if etichetta_cella_lower in chiavi_excel_target_lower:
                        chiave_standard_associata = chiavi_excel_target_lower[etichetta_cella_lower]
                        print("etichetta_cella_lower: "+ etichetta_cella_lower)
                        
                        # Il valore si trova nella colonna successiva a quella dell'etichetta
                        if etichetta_cella_lower == "indirizzo sede legale":
                            # Se l'etichetta è "indirizzo sede legale" devo prende il valrore della cella sotto 
                            # e non a destra
                            if index + 1 < len(df_excel):
                                valore_estratto = df_excel.iloc[index + 1, col_idx]
                                # Converti in stringa, gestendo valori NaN/None
                                dati_singola_azienda[chiave_standard_associata] = str(valore_estratto) if pd.notna(valore_estratto) else ""
                        elif col_idx + 1 < len(row):
                            valore_estratto = row.iloc[col_idx + 1]
                            # Converti in stringa, gestendo valori NaN/None
                            dati_singola_azienda[chiave_standard_associata] = str(valore_estratto) if pd.notna(valore_estratto) else ""
                            print(f"    Trovato '{cella_contenuto.strip()}' -> Mappato a '{chiave_standard_associata}': '{dati_singola_azienda[chiave_standard_associata]}'")
                        else:
                            print(f"    Trovato '{cella_contenuto.strip()}' (mappato a '{chiave_standard_associata}') ma non c'è una colonna successiva per il valore.")
                            dati_singola_azienda[chiave_standard_associata] = "" # Valore vuoto se non c'è colonna
        
        # Assicura che tutte le chiavi standard siano presenti nel dizionario,
        # aggiungendo quelle mancanti con un valore vuoto.
        for skey in standard_keys_final:
            if skey not in dati_singola_azienda:
                dati_singola_azienda[skey] = "" # Imposta a stringa vuota se non trovato
                print(f"    Attenzione: Chiave standard '{skey}' non trovata per {nome_societa} durante l'estrazione. Impostata a vuoto.")
        
        lista_dati_aziende.append(dati_singola_azienda)

    except FileNotFoundError:
        print(f"ERRORE: File non trovato {file_path}")
    except Exception as e:
        print(f"ERRORE durante il processamento del file {file_path}: {e}")
        import traceback
        traceback.print_exc()


# --- 3. Creazione DataFrame finale ---
if not lista_dati_aziende:
    print("\nNessun dato aziendale è stato estratto. Impossibile procedere con la creazione delle slide.")
    exit()

competitors_df = pd.DataFrame(lista_dati_aziende)
# Riordina le colonne secondo standard_keys_final se necessario, anche se DataFrame(lista_di_dict)
# tende a mantenere l'ordine delle chiavi del primo dizionario se tutte le chiavi sono presenti.
# Per sicurezza, puoi riordinare:
competitors_df = competitors_df[standard_keys_final]


print("\n--- DataFrame finale costruito ---")
print(competitors_df.head())


# --- 4. Creazione delle slide ---
print("\n--- Inizio creazione slide PowerPoint ---")
try:
    prs = Presentation("Competitors.pptx") # Assicurati che questo file esista e abbia almeno 2 slide
except Exception as e:
    print(f"Errore nell'apertura del file 'Competitors.pptx': {e}")
    print("Assicurati che il file esista nella stessa cartella dello script e non sia corrotto.")
    exit()


if len(prs.slides) < 2:
    print("ERRORE: La presentazione 'Competitors.pptx' deve avere almeno due slide.")
    print("La seconda slide (indice 1) viene usata come modello per layout e forme fisse.")
    exit()

model_slide_layout = prs.slides[1].slide_layout
# La prima slide di output sarà la seconda slide esistente (indice 1)
current_slide = prs.slides[1] 
# Se la slide modello (prs.slides[1]) ha già le intestazioni (aggiunte da add_shapes),
# non è necessario richiamare utils.add_shapes() per questa prima slide.
# Se invece è una slide "pulita" e vuoi le intestazioni, decommenta la riga sotto.
# utils.add_shapes(current_slide) 

offset_per_riga_emu = 763200 
items_per_slide = 6 
rows_written_on_current_slide = 0

print(f"Trovate {len(competitors_df)} aziende nel DataFrame da inserire nella presentazione.")

if not competitors_df.empty:
    for index, competitor_data_dict in enumerate(competitors_df.to_dict(orient='records')):
        # Crea una nuova slide se necessario (dopo la prima e ogni items_per_slide)
        if index > 0 and index % items_per_slide == 0:
            print(f"--- Creazione nuova slide per azienda indice {index} ({competitor_data_dict.get('societa', 'N/A')}) ---")
            current_slide = prs.slides.add_slide(model_slide_layout)
            rows_written_on_current_slide = 0
            utils.add_shapes(current_slide) # Aggiungi le forme fisse (intestazioni, linee) alla nuova slide

        current_y_offset = rows_written_on_current_slide * offset_per_riga_emu
        
        print(f"  Aggiungo riga {rows_written_on_current_slide + 1} alla slide corrente per: {competitor_data_dict.get('societa', 'N/A')}")
        utils.add_competitor_row(current_slide, competitor_data_dict, y_offset=current_y_offset)
        
        rows_written_on_current_slide += 1
else:
    print("Il DataFrame è vuoto. Nessuna riga da aggiungere alla presentazione.")

try:
    prs.save("Competitors_modificato_pandas.pptx")
    print("\nPresentazione salvata con successo come 'Competitors_modificato_pandas.pptx'")
except Exception as e:
    print(f"ERRORE durante il salvataggio della presentazione: {e}")

